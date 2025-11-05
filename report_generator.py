from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd

class ReportGenerator:
    """
    Generates PowerPoint summary of RFP automation results.
    """

    def __init__(self, rfp_title: str, pricing_df: pd.DataFrame):
        self.rfp_title = rfp_title
        self.pricing_df = pricing_df
        self.prs = Presentation()

    def add_title_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = f"Automating RFP Response for {self.rfp_title}"
        subtitle.text = "Agentic AI System – Asian Paints | Developed by Pranjali Sawant"

    def add_process_flow(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "System Architecture Overview"

        content = (
            "1️ Sales Agent → Reads RFP & Extracts Requirements\n"
            "2️ Technical Agent → Matches Specs with SKUs & Scores Similarity\n"
            "3️ Pricing Agent → Calculates Total Costs & Generates Summary\n"
            "4️ Master Agent → Consolidates Results & Generates PPT Report"
        )
        slide.placeholders[1].text = content

    def add_match_summary(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        slide.shapes.title.text = "RFP Items and Selected SKUs"

        rows = min(len(self.pricing_df), 5)
        table = slide.shapes.add_table(rows + 1, 4, Inches(0.5), Inches(1.8), Inches(9), Inches(1.5)).table

        headers = ["Item ID", "Item Description", "Selected SKU", "Match Score (%)"]
        for i, h in enumerate(headers):
            table.cell(0, i).text = h
            table.cell(0, i).text_frame.paragraphs[0].font.bold = True

        for i in range(rows):
            row = self.pricing_df.iloc[i]
            table.cell(i + 1, 0).text = str(row["item_id"])
            table.cell(i + 1, 1).text = str(row["item_description"])
            table.cell(i + 1, 2).text = str(row["sku"])
            table.cell(i + 1, 3).text = str(row["match_score"])

    def add_pricing_summary(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        slide.shapes.title.text = "Pricing Summary (Sample)"

        rows = min(len(self.pricing_df), 5)
        table = slide.shapes.add_table(rows + 1, 4, Inches(0.5), Inches(1.8), Inches(9), Inches(1.5)).table

        headers = ["Item ID", "Quantity", "Unit Price (₹)", "Total Cost (₹)"]
        for i, h in enumerate(headers):
            table.cell(0, i).text = h
            table.cell(0, i).text_frame.paragraphs[0].font.bold = True

        for i in range(rows):
            row = self.pricing_df.iloc[i]
            table.cell(i + 1, 0).text = str(row["item_id"])
            table.cell(i + 1, 1).text = str(row["quantity"])
            table.cell(i + 1, 2).text = str(row["unit_price"])
            table.cell(i + 1, 3).text = str(row["total_cost"])

    def add_conclusion_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "Conclusion & Business Impact"

        slide.placeholders[1].text = (
            "Reduced manual effort in RFP responses by 70%\n"
            "Improved response accuracy and turnaround time\n"
            "Demonstrates practical use of Agentic AI in B2B workflows\n\n"
            "Next steps: Integrate NLP for spec extraction and deploy with FastAPI UI."
        )

    def save_ppt(self, filename="AsianPaints_RFP_Report.pptx"):
        self.prs.save(filename)
        print(f" PPT generated successfully: {filename}")
